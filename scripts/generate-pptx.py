"""
HTML路演PPT -> PPTX 转换器
用playwright截取每一页slide，用python-pptx生成全屏背景图PPTX
"""

import os
import sys
from pathlib import Path
from lxml import etree

# 项目根目录
PROJECT_DIR = Path(__file__).resolve().parent.parent
HTML_PATH = PROJECT_DIR / "docs" / "路演ppt.html"
SLIDES_DIR = PROJECT_DIR / "docs" / "slides"
OUTPUT_PPTX = PROJECT_DIR / "docs" / "路演PPT.pptx"

SLIDES_DIR.mkdir(parents=True, exist_ok=True)

# ============ STEP 1: 用playwright截取每页slide ============

def capture_slides():
    from playwright.sync_api import sync_playwright

    html_url = HTML_PATH.as_uri()
    print(f"[1/2] 打开HTML: {html_url}")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(html_url, wait_until="networkidle")
        page.wait_for_timeout(3000)

        # 注入CSS: 强制所有动画元素立即可见, 隐藏导航UI, 暂停动画, 防溢出+放大内容
        page.add_style_tag(content="""
            /* 强制所有anim-item立即可见（跳过IntersectionObserver） */
            .anim-item {
                opacity: 1 !important;
                transform: translateY(0) !important;
                transition: none !important;
            }
            /* 隐藏导航圆点、进度条、页码计数器 */
            .nav-dots, #navDots { display: none !important; }
            .progress-bar, #progressBar { display: none !important; }
            .page-counter, #pageCounter { display: none !important; }
            /* 隐藏背景装饰粒子和光球 */
            .ambient-orb { display: none !important; }
            #particles { display: none !important; }
            /* 暂停所有CSS动画 */
            *, *::before, *::after {
                animation-play-state: paused !important;
            }
            /* 禁用smooth滚动，使scrollIntoView instant生效 */
            html {
                scroll-behavior: auto !important;
            }
            /* 放大内容区域，让PPT文字更清晰 */
            .slide-inner {
                transform: scale(1.15);
                transform-origin: center center;
            }
            /* 防止溢出 */
            .slide {
                overflow: hidden !important;
            }
            /* 确保文字不超出边界 */
            * {
                word-break: break-word;
                overflow-wrap: break-word;
            }
            /* 竞争分析表格：缩小字号防溢出 */
            .comp-table {
                font-size: 0.8rem !important;
                min-width: auto !important;
            }
            .comp-table th,
            .comp-table td {
                padding: 10px 10px !important;
                white-space: nowrap;
            }
            /* 市场同心圆SOM圈内文字缩小 */
            .market-ring.som .ring-value {
                font-size: 0.6rem !important;
                max-width: 120px !important;
                line-height: 1.2 !important;
            }
        """)
        page.wait_for_timeout(500)

        slides = page.query_selector_all("section.slide")
        slide_count = len(slides)
        print(f"   检测到 {slide_count} 页slide")

        # 先强制所有anim-item可见 + 触发计数器和柱状图动画
        page.evaluate("""
            // 强制所有anim-item立即可见
            document.querySelectorAll('.anim-item').forEach(el => {
                el.classList.add('visible');
            });
            // 强制计数器动画完成（直接设最终值）
            document.querySelectorAll('[data-count]').forEach(el => {
                el.textContent = el.dataset.count;
                el.dataset.counted = 'true';
            });
            // 强制柱状图填充到目标高度
            document.querySelectorAll('.bar-fill').forEach(el => {
                if (el.dataset.height) {
                    el.style.height = (el.dataset.height / 100 * 180) + 'px';
                    el.style.transition = 'none';
                }
            });
            var barChart = document.getElementById('barChart');
            if (barChart) barChart.dataset.animated = 'true';
        """)
        page.wait_for_timeout(500)

        screenshot_paths = []
        for i in range(slide_count):
            # 用scrollIntoView滚动到该slide（scroll容器是html，不是deck）
            page.evaluate(f"""
                const slide = document.getElementById('slide-{i}');
                if (slide) {{
                    slide.scrollIntoView({{behavior: 'instant', block: 'start'}});
                }}
            """)
            page.wait_for_timeout(800)  # 等scroll-snap定位完成

            img_path = str(SLIDES_DIR / f"slide-{i:02d}.png")
            page.screenshot(path=img_path, type="png")
            screenshot_paths.append(img_path)
            print(f"   截图 slide-{i:02d}.png ({i+1}/{slide_count})")

        browser.close()
        print(f"   所有截图已保存到 {SLIDES_DIR}")
        return screenshot_paths


# ============ STEP 2: 过渡动画 ============

def add_transition(slide, trans_type='fade'):
    """给幻灯片添加过渡动画"""
    from pptx.oxml.ns import qn

    # 移除已有transition
    for existing in slide._element.findall(qn('p:transition')):
        slide._element.remove(existing)

    transition = etree.SubElement(
        slide._element,
        qn('p:transition'),
        attrib={'spd': 'med', 'advClick': '1'}
    )

    type_map = {
        'fade': qn('p:fade'),
        'push': qn('p:push'),
        'wipe': qn('p:wipe'),
        'cover': qn('p:cover'),
        'split': qn('p:split'),
    }

    if trans_type in type_map:
        etree.SubElement(transition, type_map[trans_type])


# 过渡动画交替序列
TRANSITION_CYCLE = ['fade', 'push', 'fade', 'wipe', 'fade', 'cover']


# ============ STEP 3: 生成PPTX ============

def create_pptx(screenshot_paths):
    from pptx import Presentation
    from pptx.util import Emu, Inches

    print(f"[2/2] 生成PPTX...")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(screenshot_paths):
        if not os.path.exists(img_path):
            print(f"   警告: 截图不存在 {img_path}, 跳过")
            continue

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

        # 添加过渡动画（交替使用不同效果）
        trans = TRANSITION_CYCLE[i % len(TRANSITION_CYCLE)]
        add_transition(slide, trans)
        print(f"   添加 slide-{i:02d} 到PPTX (过渡: {trans})")

    prs.save(str(OUTPUT_PPTX))
    print(f"\n完成! PPTX已保存: {OUTPUT_PPTX}")
    print(f"   共 {len(screenshot_paths)} 页, viewport 1920x1080, 含过渡动画")


# ============ MAIN ============

if __name__ == "__main__":
    if not HTML_PATH.exists():
        print(f"错误: HTML文件不存在 {HTML_PATH}")
        sys.exit(1)

    paths = capture_slides()
    create_pptx(paths)
